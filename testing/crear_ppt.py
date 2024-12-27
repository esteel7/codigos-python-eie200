from pptx import Presentation

# Creamos nuevamente la presentación para el Módulo 2 de Microsoft Visio.

# Crear la presentación
prs_mod2 = Presentation()

# Diapositiva 1 - Título
slide_layout = prs_mod2.slide_layouts[0]
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Elaboración de Diagramas en Microsoft Visio"
subtitle.text = "Módulo 2: Diagramas Organizacionales, Redes y Efectos Visuales"

# Diapositiva 2 - Introducción
slide_layout = prs_mod2.slide_layouts[1]
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Introducción al Módulo 2"
content.text = """
En este módulo, exploraremos la creación de diagramas avanzados en Microsoft Visio, 
incluyendo diagramas organizacionales, de redes y sistemas, y cómo aplicar formatos 
y efectos visuales para mejorar la claridad y estética de los diagramas.
"""

# Diapositiva 3 - Diagramas Organizacionales
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "¿Qué es un Diagrama Organizacional?"
content.text = """
Un diagrama organizacional representa la estructura jerárquica de una organización, 
mostrando la relación entre empleados, departamentos y sus responsabilidades.
"""

# Diapositiva 4 - Uso de Diagramas Organizacionales
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "¿Por qué utilizar Diagramas Organizacionales?"
content.text = """
- Visualiza la estructura jerárquica de una empresa o equipo.
- Facilita la comunicación mostrando roles y responsabilidades.
- Ayuda en la planificación y reorganización de recursos humanos.
"""

# Diapositiva 5 - Plantillas en Visio
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Plantillas en Visio para Diagramas Organizacionales"
content.text = """
Visio ofrece plantillas prediseñadas que permiten agregar y personalizar formas de 
manera rápida. Además, puedes vincular datos externos, como hojas de Excel, para 
actualizar automáticamente el diagrama.
"""

# Diapositiva 6 - Diagramas de Redes y Sistemas
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "¿Qué son los Diagramas de Redes y Sistemas?"
content.text = """
Los diagramas de redes y sistemas permiten visualizar la infraestructura de TI, 
mostrando cómo los dispositivos y servidores están conectados.
"""

# Diapositiva 7 - Usos de Diagramas de Redes
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "¿Por qué utilizar Diagramas de Redes?"
content.text = """
- Mapear la infraestructura de TI.
- Visualizar conexiones entre routers, switches, servidores y usuarios.
- Facilitar la planificación y solución de problemas.
"""

# Diapositiva 8 - Plantillas en Visio para Redes
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Plantillas en Visio para Diagramas de Redes"
content.text = """
Visio ofrece plantillas especializadas que incluyen formas para routers, switches, 
servidores y otros dispositivos de red.
"""

# Diapositiva 9 - Efectos Visuales
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Efectos Visuales en Visio"
content.text = """
Los efectos visuales ayudan a resaltar elementos clave en un diagrama. Se pueden 
aplicar colores, sombras, líneas gruesas y efectos de 3D para mejorar la claridad 
y estética.
"""

# Diapositiva 10 - Temas y Estilos
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Aplicación de Temas y Estilos"
content.text = """
Visio ofrece temas y estilos predefinidos que puedes aplicar a todo el diagrama. 
Esto incluye combinaciones de colores, estilos de línea y tipos de fuente.
"""

# Diapositiva 11 - Personalización
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Personalización de Formatos"
content.text = """
Se pueden personalizar colores, grosores de línea y estilos de fuente para 
adaptar el diagrama a la identidad visual de la organización.
"""

# Diapositiva 12 - Efectos 3D
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Aplicación de Efectos 3D y Sombra"
content.text = """
Visio permite aplicar efectos de 3D y sombra a las formas para agregar profundidad 
y resaltar jerarquías importantes.
"""

# Diapositiva 13 - Ejemplo de Diagrama de Red
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Ejemplo: Diagrama de Red"
content.text = "A continuación, un ejemplo de diagrama de red que representa una pequeña infraestructura de TI."

# Diapositiva 14 - Ejemplo Visual
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Ejemplo: Diagrama con Efectos Visuales"
content.text = "Este diagrama utiliza colores y efectos visuales para mejorar la claridad y estética."

# Diapositiva 15 - Ejercicio Práctico
slide = prs_mod2.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Ejercicio Práctico"
content.text = """
Crea un diagrama organizacional y un diagrama de red utilizando plantillas de Visio. 
Aplica estilos y efectos visuales para destacar elementos clave.
"""

# Guardar la presentación
prs_mod2.save("G:\Mi unidad\Docencia\PUCV\Programación numérica\Python Codes\pptx\Presentacion_Microsoft_Visio_Modulo_2.pptx")
