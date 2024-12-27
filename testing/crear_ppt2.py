from pptx import Presentation
from pptx.util import Inches

# Crear una nueva presentación
prs = Presentation()

# 1. Portada del Módulo 2
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Visualizaciones Avanzadas en Power BI"
subtitle.text = "Creación y Personalización de Gráficos, Filtros y Paneles"

# 2. Objetivos del Módulo
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Objetivos del Módulo"
content = slide.shapes.placeholders[1].text = (
    "• Introducir diferentes tipos de visualizaciones.\n"
    "• Aprender a personalizar gráficos.\n"
    "• Aplicar segmentaciones y filtros avanzados.\n"
    "• Crear paneles interactivos y gráficos jerárquicos."
)

# 3. Tipos de Visualizaciones Avanzadas
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Tipos de Visualizaciones Avanzadas"
content = slide.shapes.placeholders[1].text = (
    "• Gráficos de barras.\n"
    "• Gráficos circulares.\n"
    "• Gráficos de dispersión.\n"
    "• Mapas geográficos.\n"
    "• Diagramas de árbol."
)

# 4. Ejemplo 1: Gráficos de Barras
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Ejemplo 1: Gráficos de Barras y Personalización"
content = slide.shapes.placeholders[1].text = (
    "1. Importa los datos de ventas desde Excel.\n"
    "2. Selecciona 'Gráfico de Barras Apiladas'.\n"
    "3. Eje X: Coloca la columna 'Categoría'.\n"
    "4. Valores: Coloca 'Ventas 2023' y 'Ventas 2022'.\n"
    "5. Personaliza el gráfico cambiando colores y agregando etiquetas de datos.\n"
    "6. Filtra por categoría para una visualización interactiva."
)

# 5. Ejemplo 2: Gráficos de Dispersión
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Ejemplo 2: Gráficos de Dispersión"
content = slide.shapes.placeholders[1].text = (
    "1. Importa los datos de productos desde Excel.\n"
    "2. Selecciona 'Gráfico de Dispersión'.\n"
    "3. Eje X: Precio.\n"
    "4. Eje Y: Unidades Vendidas.\n"
    "5. Personaliza los puntos con colores y tamaños según el margen de ganancias.\n"
    "6. Filtra por productos para un análisis más detallado."
)

# 6. Filtros y Segmentaciones
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Uso de Filtros y Segmentaciones"
content = slide.shapes.placeholders[1].text = (
    "• Agrega filtros avanzados para controlar la visualización de los datos.\n"
    "• Ejemplo: Filtrar por año o categoría de producto.\n"
    "• Usa segmentaciones para hacer visualizaciones interactivas más flexibles."
)

# 7. Paneles y Cuadros de Mando
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Creación de Paneles y Cuadros de Mando"
content = slide.shapes.placeholders[1].text = (
    "• Un panel en Power BI combina gráficos, tablas y otros elementos visuales.\n"
    "• Ejemplo: Un panel de ventas que muestra KPIs, gráficos de tendencia y mapas.\n"
    "• Personaliza el diseño y asegúrate de que sea intuitivo para los usuarios."
)

# 8. Mapas Geográficos y Datos Espaciales
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Mapas Geográficos y Datos Espaciales"
content = slide.shapes.placeholders[1].text = (
    "• Visualiza datos basados en la ubicación geográfica.\n"
    "• Ejemplo: Un mapa que muestra ventas por región.\n"
    "• Usa datos espaciales para crear mapas interactivos."
)

# 9. Gráficos de Líneas de Tiempo
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Gráficos de Líneas de Tiempo"
content = slide.shapes.placeholders[1].text = (
    "• Ideal para analizar tendencias en series temporales.\n"
    "• Ejemplo: Gráfico de ventas mensuales a lo largo del año.\n"
    "• Agrega herramientas de predicción para análisis de tendencias futuras."
)

# 10. Herramientas de Análisis Predictivo
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Herramientas de Análisis Predictivo"
content = slide.shapes.placeholders[1].text = (
    "• Power BI ofrece herramientas para análisis predictivo como el pronóstico.\n"
    "• Ejemplo: Aplicar el pronóstico a un gráfico de ventas para predecir valores futuros.\n"
    "• Útil para planificar estrategias a largo plazo basadas en tendencias históricas."
)

# 11. Conclusión del Módulo
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Conclusión del Módulo"
content = slide.shapes.placeholders[1].text = (
    "• En este módulo, aprendimos a crear visualizaciones avanzadas, paneles y filtros interactivos.\n"
    "• Explora las herramientas de personalización y predicción para mejorar la toma de decisiones.\n"
    "• Experimenta con diferentes tipos de gráficos para obtener el máximo valor de tus datos."
)

# Guardar la presentación
#file_path = "/ppt/PowerBI_Modulo2_Visualizaciones_Avanzadas.pptx"
#prs.save(file_path)

#file_path

# Guardar la presentación
prs.save("G:\Mi unidad\Docencia\PUCV\Programación numérica\Python Codes\pptx\Presentacion_Power_BI_Modulo_2.pptx")